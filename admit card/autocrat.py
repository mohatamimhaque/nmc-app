#!/usr/bin/env python3
"""
autocrat.py

FastAPI admit card document generator.
- Listens locally on port 8000.
- Accepts POST /generate-admit-card with registration parameters.
- Validates the X-API-Key header.
- Replaces placeholders in template.pptx, generates a QR code from serial, and converts to PDF.
- Text placeholders are replaced using Poppins Medium font at 20px size.
- Cleans up all temporary files immediately.
"""

import os
import re
import sys
import time
import uuid
import threading
import qrcode
from typing import Dict, Any
from fastapi import FastAPI, HTTPException, Header, status
from fastapi.responses import Response
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt

# Load environment variables from .env.local in parent directory
env_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".env.local"))
if os.path.exists(env_path):
    load_dotenv(dotenv_path=env_path)

secret_key = os.getenv("ADMIT_CARD_SECRET_KEY")

app = FastAPI(title="NMC Admit Card Generator API")

class PPTXToPDFConverter:
    """Manages PowerPoint COM conversion locally inside thread-safe context"""
    def convert(self, pptx_path, pdf_path):
        import pythoncom
        import win32com.client
        
        # Initialize COM library for this thread
        pythoncom.CoInitialize()
        powerpoint = None
        try:
            # Dispatch uses an existing PowerPoint instance or opens a new one
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
            deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 is ppSaveAsPDF
            deck.Close()
            return True
        except Exception as e:
            print(f"    [WARN] PowerPoint conversion failed: {str(e)[:150]}")
            return False
        finally:
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except:
                    pass
            # Uninitialize COM for this thread
            pythoncom.CoUninitialize()

# Instantiate converter & lock
converter = PPTXToPDFConverter()
powerpoint_lock = threading.Lock()


def normalize_key(k):
    """Remove spaces and make key consistent for placeholder matching"""
    return str(k).replace(" ", "")

def replace_placeholders_in_text(text, mapping):
    if text is None:
        return text
    def repl(m):
        key = m.group(1).strip()
        return str(mapping.get(key, m.group(0)))
    return re.sub(r"\{\{\s*([^}]+)\s*\}\}", repl, text)

def generate_filled_pptx(template_path, row_map, out_path):
    """Generate a filled PPTX presentation replacing placeholders including {{qr}}"""
    # Extract Serial column value case-insensitively from row_map
    serial_val = ""
    for k, v in row_map.items():
        if k.lower() == 'serial':
            serial_val = str(v).strip()
            break
            
    temp_qr_path = None
    if serial_val:
        qr = qrcode.QRCode(version=1, box_size=10, border=1)
        qr.add_data(serial_val)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        temp_qr_path = out_path + "_qr_temp.png"
        img.save(temp_qr_path)
        
    prs = Presentation(template_path)
    for slide in prs.slides:
        shapes_to_delete = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Check if this shape is the QR placeholder
            text_content = shape.text_frame.text
            if re.search(r"\{\{\s*qr\s*\}\}", text_content, re.IGNORECASE):
                if temp_qr_path and os.path.exists(temp_qr_path):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Make it a square of exactly 250px (1 pixel = 9525 EMUs in Office documents)
                    square_size = 400 * 9525
                    new_left = left + (width - square_size) // 2
                    new_top = top + (height - square_size) // 2
                    
                    # Add picture at centered square position
                    slide.shapes.add_picture(temp_qr_path, new_left, new_top, square_size, square_size)
                    shapes_to_delete.append(shape)
                continue
                
            # Replace other placeholders
            for para in shape.text_frame.paragraphs:
                full = "".join(run.text for run in para.runs)
                if "{{" in full:
                    new = replace_placeholders_in_text(full, row_map)
                    for i in range(len(para.runs)-1, -1, -1):
                        try:
                            para.runs[i].text = ""
                        except Exception:
                            pass
                    
                    if len(para.runs) > 0:
                        run = para.runs[0]
                        run.text = new
                        run.font.name = 'Poppins Medium'
                        run.font.size = Pt(20)
                    else:
                        run = para.add_run()
                        run.text = new
                        run.font.name = 'Poppins Medium'
                        run.font.size = Pt(20)
                        
        # Delete placeholder shapes
        for shape in shapes_to_delete:
            try:
                shape._element.getparent().remove(shape._element)
            except Exception as e:
                # Fallback in case of removal error
                try:
                    shape.text_frame.text = ""
                except:
                    pass
                    
    prs.save(out_path)
    
    # Clean up temp QR code image
    if temp_qr_path and os.path.exists(temp_qr_path):
        try:
            os.remove(temp_qr_path)
        except Exception:
            pass

def safe_remove(file_path, retries=5, delay=0.2):
    """Attempt to delete a file multiple times in case it is locked by another process"""
    for i in range(retries):
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
            return True
        except Exception:
            time.sleep(delay)
    return False

# Fuzzy aliases mapping
MAPPING_RULES = {
    "Name": ["fullname", "name", "participant_name", "participantname", "user_name", "username", "participant"],
    "Institution": ["institution", "institution_name", "institutionname", "school", "college", "university"],
    "Phone_Number": ["phone_number", "phonenumber", "phone", "mobile", "mobile_number", "mobilenumber", "contact", "contact_number"],
    "Participation_Level": ["level", "participation_level", "participationlevel", "category", "class", "grade"],
    "Select_Event": ["event", "event_name", "eventname", "select_event", "selectevent"],
    "Serial": ["serial", "registration_id", "registrationid", "reg_id", "regid", "serial_number", "serialnumber", "id"]
}

def clean_key(k: str) -> str:
    return re.sub(r'[^a-z0-9]', '', str(k).lower())

def remap_parameters(input_data: Dict[str, Any]) -> Dict[str, Any]:
    remapped = {}
    cleaned_input = {clean_key(k): v for k, v in input_data.items()}
    
    cleaned_db_cols = {
        "Name": clean_key("full_name"),
        "Institution": clean_key("institution"),
        "Phone_Number": clean_key("phone_number"),
        "Participation_Level": clean_key("level"),
        "Select_Event": clean_key("event"),
        "Serial": clean_key("serial")
    }
    
    for target_key, aliases in MAPPING_RULES.items():
        # 1st Priority: Exact key (normalized)
        norm_target = clean_key(target_key)
        if norm_target in cleaned_input:
            remapped[target_key] = cleaned_input[norm_target]
            continue
            
        # 2nd Priority: Database column name (normalized)
        db_col = cleaned_db_cols[target_key]
        if db_col in cleaned_input:
            remapped[target_key] = cleaned_input[db_col]
            continue
            
        # 3rd Priority: Aliases
        matched_val = None
        for alias in aliases:
            norm_alias = clean_key(alias)
            if norm_alias in cleaned_input:
                matched_val = cleaned_input[norm_alias]
                break
        if matched_val is not None:
            remapped[target_key] = matched_val
            continue
            
    # "if not match when send fix it" - ensure all target parameters are present
    for target_key in MAPPING_RULES.keys():
        if target_key not in remapped:
            remapped[target_key] = ""
            
    return remapped

@app.post("/generate-admit-card")
def generate_admit_card(
    payload: Dict[str, Any],
    x_api_key: str = Header(None, alias="X-API-Key")
):
    if not secret_key:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="ADMIT_CARD_SECRET_KEY is not configured on the server."
        )
    if x_api_key != secret_key:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Unauthorized access."
        )
        
    # Remap parameters based on template rules
    remapped_params = remap_parameters(payload)
    serial = str(remapped_params.get("Serial", "")).strip()
    
    if not serial:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Serial parameter is required to generate the admit card."
        )
        
    safe_serial = re.sub(r'[\\/*?:"<>|]', "_", serial)
    pptx_dir = os.path.dirname(os.path.abspath(__file__))
    temp_dir = os.path.join(pptx_dir, "temp")
    os.makedirs(temp_dir, exist_ok=True)
    
    temp_pptx_path = os.path.join(temp_dir, f"{safe_serial}_{uuid.uuid4().hex[:8]}_temp.pptx")
    temp_pdf_path = os.path.join(temp_dir, f"{safe_serial}_{uuid.uuid4().hex[:8]}_temp.pdf")
    template_pptx = os.path.join(pptx_dir, "template.pptx")
    
    if not os.path.exists(template_pptx):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Template PowerPoint file 'template.pptx' not found."
        )
        
    with powerpoint_lock:
        try:
            # 1. Fill PowerPoint template
            generate_filled_pptx(template_pptx, remapped_params, temp_pptx_path)
            
            # 2. Convert to PDF
            success = converter.convert(temp_pptx_path, temp_pdf_path)
            if not success:
                raise HTTPException(
                    status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                    detail="Failed to convert PowerPoint to PDF."
                )
        except Exception as e:
            safe_remove(temp_pptx_path)
            safe_remove(temp_pdf_path)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error generating PDF: {str(e)}"
            )
        finally:
            safe_remove(temp_pptx_path)
            
    try:
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to read generated PDF file: {str(e)}"
        )
    finally:
        safe_remove(temp_pdf_path)
        
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="admit_{safe_serial}.pdf"'
        }
    )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("autocrat:app", host="127.0.0.1", port=8080, reload=True)